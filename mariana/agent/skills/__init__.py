"""Skill library — domain-specific system-prompt snippets injected into the
planner based on lightweight intent classification of the user's goal.

Each skill is a tuple of (name, keywords, guidance).  The planner scores the
goal against the keyword set of each skill and injects the matching guidance
into the system prompt *after* the base prompt and tool manifest.  Non-matching
skills cost nothing.

Principles:
- Skills are additive.  They never replace the base prompt.
- Skills are LLM-agnostic.  No model-specific tokens.
- Skills encode best practices, pitfalls, and preferred tooling — not code.
- A goal may match multiple skills; top-2 by score are injected.
"""
from __future__ import annotations

import re
from dataclasses import dataclass


@dataclass(frozen=True)
class Skill:
    name: str
    keywords: tuple[str, ...]
    guidance: str


# ---------------------------------------------------------------------------
# Skill definitions
# ---------------------------------------------------------------------------


CODING = Skill(
    name="coding",
    keywords=(
        "code", "coding", "program", "script", "implement", "function",
        "class", "module", "library", "api", "cli", "tool", "debug", "refactor",
        "unit test", "pytest", "lint", "typecheck", "compile", "build",
        "python", "rust", "go", "typescript", "javascript", "c++", "java",
    ),
    guidance="""SKILL: Coding (production-grade software engineering)

When the goal involves writing, debugging, or refactoring code:
- Prefer strongly-typed languages with explicit error handling.  When writing
  Python, include type hints on every function signature and at least one
  self-test via `assert` or `pytest`.  When writing Rust/Go/TS/C++, use
  idiomatic error types and avoid panics / unchecked unwraps.
- Start by listing the expected inputs, outputs, edge cases, and failure modes
  in a `think` step BEFORE writing code.
- After writing code, ALWAYS run it (code_exec / bash_exec) with representative
  inputs and verify observable behaviour, not just exit code.
- Cover at minimum: happy path, empty input, malformed input, boundary values.
- If a test fails, inspect stderr in detail before attempting a fix — do not
  guess.  The fix loop has a hard cap of 5 attempts per step.
- Keep files small and focused; prefer multiple files over one giant file when
  the module count exceeds ~3 logical groupings.
""",
)


WEB_DEV = Skill(
    name="web_dev",
    keywords=(
        "website", "web app", "webapp", "frontend", "backend", "react",
        "next.js", "nextjs", "vue", "svelte", "tailwind", "html", "css",
        "http server", "rest api", "graphql", "express", "fastapi", "flask",
        "vite", "webpack", "deploy", "vercel", "netlify", "dockerfile",
    ),
    guidance="""SKILL: Web Development

When the goal involves building web apps, HTTP APIs, or frontends:
- For static sites, write an index.html that renders correctly when opened
  directly from the filesystem.  Inline CSS/JS unless the task demands a
  multi-file structure.
- For APIs, always define: route table, request/response schemas, error codes,
  auth model.  Pick FastAPI for Python, Axum for Rust, Fastify/Hono for Node.
- Wire a health check (`/api/health` returning `{status:"ok"}`) before any
  business logic so you can ping it from bash_exec.
- For frontend frameworks, run `npm install` and `npm run build` — do not
  ship unbuilt source.  Confirm the build succeeds.
- Never hardcode secrets.  Use placeholder env vars and document them.
""",
)


DATA_ANALYSIS = Skill(
    name="data_analysis",
    keywords=(
        "data", "dataset", "csv", "parquet", "dataframe", "pandas", "polars",
        "analyze", "analysis", "statistics", "statistical", "aggregate",
        "groupby", "join", "sql", "query", "visualize", "plot", "chart",
        "histogram", "scatter", "correlation", "regression", "anova",
    ),
    guidance="""SKILL: Data Analysis

When the goal involves exploring, transforming, or summarising data:
- Load the data once, inspect shape / dtypes / nulls / ranges BEFORE any
  aggregation.  Log these as a think step.
- Prefer polars > pandas for anything >100k rows.  Use lazy frames where
  available.
- For stats claims, compute the actual statistic (mean, median, std, p-value,
  effect size) and print it.  Never describe a trend without a number.
- Save intermediate artefacts (cleaned_data.parquet, summary.csv) to the
  workspace so downstream steps can reuse them.
- For visualisations, write charts to /workspace/<user>/<name>.png and include
  the filename in the deliver step so the user can view them.
- Watch for: unit confusion, timezone handling, missing-value semantics,
  survivorship bias.
""",
)


BROWSER_AUTOMATION = Skill(
    name="browser_automation",
    keywords=(
        "browser", "playwright", "selenium", "scrape", "crawl", "extract",
        "navigate", "click", "form", "login", "website", "web page",
        "screenshot", "pdf", "headless", "url", "http", "html",
    ),
    guidance="""SKILL: Browser Automation

When the goal requires fetching content from a live website or driving a
browser:
- Prefer `browser_fetch` for simple GET-then-read scenarios.  Use the full
  Playwright tool only when the page depends on JavaScript rendering or
  requires interaction.
- Set explicit timeouts (20-45s for network-heavy pages).  Use
  wait_for="domcontentloaded" unless you need full load.
- Always save raw HTML / screenshots to the workspace as an artefact before
  attempting parsing — it lets you inspect failures.
- Parse with a real parser (BeautifulSoup, lxml, cheerio) — never regex HTML.
- Respect robots.txt and rate limits.  One request at a time unless the goal
  explicitly requires crawling.
- If a site blocks headless browsers, try a plain HTTP fetch with realistic
  headers first before escalating.
""",
)


DEVOPS = Skill(
    name="devops",
    keywords=(
        "deploy", "docker", "container", "kubernetes", "k8s", "compose",
        "nginx", "systemd", "ci", "cd", "pipeline", "terraform", "ansible",
        "server", "linux", "ssh", "bash", "shell", "cron", "monitor",
    ),
    guidance="""SKILL: DevOps / System Administration

When the goal involves deployment, container, or server operations:
- Start with a dry run.  Never run a destructive command (rm -rf, docker
  system prune, kubectl delete) without dumping the target state first.
- For Docker / compose files, validate with `docker compose config` before
  `up`.  Always add health checks and restart policies.
- For bash scripts, `set -euo pipefail` at the top.  Quote variables.
- Capture logs.  If a service fails to start, run `docker logs <name>` or
  `journalctl -u <name>` and include the tail in your analysis.
- Idempotency: assume the script will run twice.  Use `mkdir -p`, guard
  `docker run --rm` with `docker ps`, and prefer declarative tools.
- NEVER print secrets to stdout.  Redact tokens in any logs you emit.
""",
)


ML = Skill(
    name="ml",
    keywords=(
        "model", "machine learning", "ml", "deep learning", "neural",
        "train", "training", "fine-tune", "finetune", "inference", "predict",
        "classification", "regression", "embedding", "transformer", "pytorch",
        "tensorflow", "scikit", "sklearn", "cnn", "rnn", "lstm", "llm",
    ),
    guidance="""SKILL: Machine Learning / Modeling

When the goal involves training or evaluating a model:
- Split data into train/val/test BEFORE any feature engineering.  Never touch
  the test set during iteration.
- Report the full metric set: accuracy/F1/AUC for classification; MAE/RMSE/R^2
  for regression.  Include a trivial baseline (mean predictor, majority class)
  for context.
- Seed all RNGs (numpy, torch, random) and record the seed.
- For deep learning, print model parameter count, training loss per epoch, and
  validation metric per epoch.
- For inference/deployment, measure latency (p50/p95) on a realistic payload.
- Guard against data leakage: confirm no target column survived into features,
  no future data leaked into past, no duplicate rows straddle splits.
""",
)


CONTENT = Skill(
    name="content",
    keywords=(
        "write", "article", "blog", "post", "essay", "report", "document",
        "memo", "summary", "summarise", "summarize", "translate",
        "email", "letter", "draft", "copy", "caption", "tweet",
    ),
    guidance="""SKILL: Content & Writing

When the goal is producing written prose:
- Identify the audience, length target, and tone from the goal.  Ask the user
  for clarification only if these are genuinely ambiguous AND the task cannot
  proceed without them.
- Structure longer pieces with a clear opening hook, numbered or named
  sections, and a concrete takeaway/close.
- Cite every factual claim with a URL in markdown link form.  Never invent
  citations.  If a source can't be verified, flag it in-line with `[unsourced]`.
- Save the piece to a markdown file in the workspace and include the path in
  deliver.  Large deliverables may also produce a PDF via pandoc or a docx
  via python-docx.
""",
)


RESEARCH = Skill(
    name="research",
    keywords=(
        "research", "investigate", "find out", "compare", "evaluate",
        "survey", "review", "analyse", "analyze", "study", "report",
        "market", "competitor", "landscape", "trend", "whitepaper",
    ),
    guidance="""SKILL: Research & Investigation

When the goal is to gather information from external sources and synthesize:
- Plan the question decomposition in a think step: main question → sub-
  questions → sources that would answer each.
- Use web_search / browser_fetch for primary sources.  Prefer official
  documentation, SEC filings, arxiv, vendor blogs over aggregator summaries.
- Keep a sources.md file in the workspace where you log every URL visited with
  a 1-line takeaway.  This is your evidence trail.
- Resolve conflicting data by checking the most authoritative primary source.
  Never average two contradictory numbers.
- Deliverable should distinguish: established facts, estimates, opinions, and
  unknowns.  Label each section accordingly.
""",
)


FINANCE = Skill(
    name="finance",
    keywords=(
        "stock", "equity", "bond", "options", "futures", "crypto", "forex",
        "yield", "volatility", "sharpe", "alpha", "backtest", "portfolio",
        "hedge", "arbitrage", "var", "cvar", "risk", "financial", "earnings",
        "valuation", "dcf", "macro", "interest rate", "fed", "inflation",
    ),
    guidance="""SKILL: Finance / Quantitative Research

When the goal involves markets, portfolios, trading, or financial modelling:
- Be explicit about the data source, frequency, and lookback window.  Never
  compute a return without specifying the denominator (t-1 close, prior day,
  etc.).
- For backtests: declare the universe, rebalance frequency, transaction cost
  model, position sizing rule, and out-of-sample split up front.  Report
  Sharpe, Sortino, max drawdown, calmar, hit rate, average win/loss.
- Avoid lookahead bias — features at time t must only use data available at
  time t (or earlier).
- For risk metrics, state the confidence level (e.g. 95% 1-day VaR) and the
  method (historical / parametric / Monte Carlo).
- Never recommend a specific security or position size as investment advice;
  frame results as analysis of historical data.
""",
)


SECURITY = Skill(
    name="security",
    keywords=(
        "security", "vulnerability", "exploit", "sandbox", "sanitize",
        "injection", "xss", "csrf", "ssrf", "privilege", "authz", "authn",
        "authenticate", "authorize", "encryption", "tls", "ssl", "hash",
        "password", "token", "jwt", "secret", "audit",
    ),
    guidance="""SKILL: Security Engineering

When the goal touches on auth, data handling, or exposed endpoints:
- Threat-model first.  Write out: who the attacker is, what they can do, what
  they want to steal.  Then design defences for each.
- Never trust user input.  Validate type, length, and allowed set before use.
  Parameterise every SQL query.  Sanitise every filesystem path (no ../, no
  absolute paths, no symlinks pointing outside the sandbox).
- Use constant-time comparison for secrets / tokens / MAC values.
- Prefer well-audited libraries over bespoke crypto.  Never implement your own
  primitive.
- Log auth events (success + failure) with enough context to audit, but never
  log the secret itself.
- When reviewing code, flag every place where a privilege boundary is crossed
  and every place user input reaches a sink (exec, eval, SQL, file open, shell).
""",
)


# ---------------------------------------------------------------------------
# v3 skills
# ---------------------------------------------------------------------------


SOCIAL_OPS = Skill(
    name="social_ops",
    keywords=(
        "social media", "twitter", "tweet", "x.com", "instagram", "tiktok",
        "linkedin", "facebook", "threads", "post", "caption", "hashtag",
        "schedule", "calendar", "content plan", "campaign", "influencer",
        "engagement", "follower", "reach", "impressions",
    ),
    guidance="""SKILL: Social-media Operations

When the goal involves drafting, scheduling, or analysing social-media
content:
- For multi-platform output, produce ONE content brief first (angle,
  audience, CTA, hook) and then derive per-platform variants (character
  limits, tone shift, aspect ratio).
- Character limits: X = 280, LinkedIn = 3000, Instagram caption = 2200,
  TikTok caption = 2200, YouTube short description = 5000.  Hard-validate
  length before delivery.
- Hashtags: 2-3 on LinkedIn, 3-5 on X, 8-15 on Instagram/TikTok.  No
  banned or shadow-banned tags.  Never fabricate trending tags — pull from
  a current source when relevance matters.
- Image/video posts: prefer generate_image / generate_video with an
  explicit aspect ratio per platform (1:1, 4:5, 9:16).  Save to
  assets/social/<platform>/.
- Output a single schedule.csv with columns: platform, post_time (ISO),
  asset_path, body, hashtags, cta — makes copy-paste into buffer/later
  trivial.
- Never post on a user's behalf without explicit per-platform consent.
""",
)


SALES_OPS = Skill(
    name="sales_ops",
    keywords=(
        "sales", "prospect", "lead", "outreach", "cold email", "sequence",
        "crm", "pipeline", "deal", "quota", "enrichment", "icp", "persona",
        "account", "abm", "forecast", "win rate", "conversion",
    ),
    guidance="""SKILL: Sales Operations

When the goal involves prospect research, outreach drafting, or pipeline
analysis:
- Start by confirming the ICP (industry, company size, role, geography)
  from the goal.  If any dimension is unspecified and assumed, flag it.
- Enrichment: for each prospect produce a structured row with name, title,
  company, source URL, 1-line rationale, suggested hook.  Never invent an
  email — leave blank if not verifiable.
- Outreach copy: open with a specific observation (not 'I noticed your
  company'), one value claim with evidence, one soft CTA.  Max 90 words.
- Respect CAN-SPAM / GDPR: include unsubscribe language if producing actual
  send-ready copy, and note jurisdiction constraints in the deliverable.
- Never submit sends or edit a live CRM without an explicit approval step.
""",
)


RECONCILIATION = Skill(
    name="reconciliation",
    keywords=(
        "reconcile", "reconciliation", "invoice", "billing", "payment",
        "ap", "ar", "accounts payable", "accounts receivable", "ledger",
        "journal", "bank statement", "expense", "receipt", "match",
        "discrepancy", "variance", "audit trail", "chart of accounts",
    ),
    guidance="""SKILL: Reconciliation / Back-office finance

When the goal involves matching transactions across systems (bank <-> ledger,
invoice <-> payment, expense <-> receipt):
- Load every input independently, print row counts + date ranges, and
  confirm they overlap BEFORE any matching.
- Match on (amount, date within tolerance, counterparty).  Tolerance
  defaults: amount = exact, date = +/- 3 business days.  Document the
  tolerances.
- Output three tables: matched, unmatched-in-A, unmatched-in-B.  Include
  a reason code on every unmatched row (duplicate, out-of-period,
  amount-off, counterparty-missing).
- Never silently drop rows.  Every input row must be accounted for in
  exactly one output table.
- Deliverable format: XLSX with three sheets (Matched, Unmatched-Left,
  Unmatched-Right) and a Summary sheet with totals by reason code.
- Record the raw inputs' SHA256 on the Summary sheet so the audit trail
  is cryptographically linked to the reconciliation.
""",
)


FINANCE_MODELING = Skill(
    name="finance_modeling",
    keywords=(
        "financial model", "three-statement", "3-statement", "dcf",
        "valuation", "wacc", "terminal value", "lbo", "merger model",
        "accretion", "dilution", "pro forma", "cap table", "waterfall",
        "scenarios", "sensitivity", "monte carlo", "operating model",
    ),
    guidance="""SKILL: Financial Modelling

When the goal involves building a financial model:
- Structure: separate Inputs, Assumptions, Calcs, and Outputs sheets.
  Colour code: blue = hard input, black = formula, green = link to other
  sheet.
- Use real Excel formulas (SUM, INDEX/MATCH, XLOOKUP, IFERROR) so the
  user can tweak assumptions.  Never paste calculated values.  Define
  named ranges for key drivers.
- Three-statement models: balance sheet MUST balance at every period;
  add a diagnostic row `Assets - (Liab + Equity)` that should be zero,
  and highlight if non-zero.
- Scenario manager: build Base/Bull/Bear with a single cell switch.  Use
  CHOOSE() or INDEX() driven by a dropdown data-validation cell.
- Sensitivity tables: use Excel data tables (row/col) not hard-coded
  grids.  State the two driving variables in the header.
- End deliverable should include a README cell on the first sheet
  describing file structure, assumption sources, and model date.
""",
)


DOCUMENT_REVIEW = Skill(
    name="document_review",
    keywords=(
        "review", "proofread", "redline", "fact-check", "verify", "audit",
        "check", "qa", "quality assurance", "edit", "critique", "feedback",
        "contract review", "nda", "compliance check", "consistency",
    ),
    guidance="""SKILL: Document Review

When the goal is to review or fact-check an existing document:
- Read the document end-to-end ONCE before offering any edits.  Capture
  the document's claim structure first.
- Produce findings in four buckets: factual errors, internal
  inconsistencies (numbers that don't add up, names that change), style
  issues, unclear passages.  Bucket every finding; never mix.
- For every factual claim: attempt verification via web_search /
  browser_fetch.  Cite the primary source and label the claim as
  verified, contradicted, or unverifiable.
- Numerical checks: re-compute every sum / percentage / ratio in the
  document.  Report the computed value alongside the stated value.
- Deliverable: a markdown findings report with one row per finding
  (location, severity, issue, suggested fix, evidence).  Never silently
  rewrite the document — suggestions go next to the original, not over
  it.
""",
)


DELIVERABLE_QUALITY = Skill(
    name="deliverable_quality",
    keywords=(
        "report", "deliverable", "document", "pdf", "docx", "pptx",
        "presentation", "slides", "deck", "spreadsheet", "excel", "xlsx",
        "workbook", "dashboard", "memo", "brief", "whitepaper", "write-up",
        "summary", "summarize", "write up", "write a", "build a",
    ),
    guidance="""SKILL: Deliverable quality (Mariana house style)

Deliverables are what the user sees.  Polish here determines perception
of the entire system.  House style: Steve Jobs / FT / Stripe \u2014 precise,
low-chrome, high-signal.

## Format selection
- Default = Markdown (.md).  The UI renders MD inline and offers a PDF
  download button, so MD covers 80% of cases at zero risk.
- Use XLSX when the user says \"spreadsheet\", \"excel\", \"model\",
  \"reconcile\", or the output is rows\u00d7columns of numbers.
- Use PPTX when the user says \"slides\", \"deck\", \"presentation\", or
  needs a live-pitch artefact.
- Use DOCX when the user says \"word\", \"redline\", or the output will be
  co-edited in Word.
- Use PDF only when the user explicitly asks, or for locked-artefact
  distribution (contracts, final memos).  Otherwise ship MD and let the
  UI produce the PDF.

## Every deliverable MUST include
1. A title and one-line summary at the top.
2. A clearly labelled \"As of\" date.
3. A \"Sources\" section listing every URL consulted (with access date).
  Also write `sources.json` alongside the deliverable with the same
  list in structured form \u2014 one object per source with keys
  {url, title, accessed_at}.
4. A \"Methodology\" block of 2-5 bullets explaining how the answer was
  produced.  Name the tools used (web_search, browser_fetch, etc.).
5. No placeholders, no TBDs, no Lorem Ipsum.  Ship only finished work.

## Markdown deliverables
- Use `#` for title, `##` for sections, `###` sparingly.
- Tables are GitHub-flavoured.  Right-align numeric columns with `---:`.
- Inline citations as `[label](url)` \u2014 never bare URLs.
- Final file at `/workspace/<name>.md`.

## XLSX deliverables (openpyxl)
- Always write real formulas, never pasted numbers.  Wrong:
  `ws['C2'] = 12000`.  Right: `ws['C2'] = '=A2*B2'`.
- Sheets: Inputs (blue fill), Assumptions, Calcs, Outputs.  Use
  `PatternFill(start_color='DCE6F1', fill_type='solid')` for inputs.
- Number format: `'#,##0'` for currency, `'0.0%'` for percentages,
  `'0.00'` for ratios.
- Every sheet: freeze row 1 (`ws.freeze_panes = 'A2'`), set column
  widths to accommodate headers.
- Recipe:
  ```python
  from openpyxl import Workbook
  from openpyxl.styles import Font, PatternFill, Alignment
  wb = Workbook(); ws = wb.active; ws.title = 'Inputs'
  ws['A1'] = 'Driver'; ws['B1'] = 'Value'
  for c in ('A1','B1'): ws[c].font = Font(bold=True)
  ws.freeze_panes = 'A2'
  ws.column_dimensions['A'].width = 24
  wb.save('/workspace/out.xlsx')
  ```

## PPTX deliverables (python-pptx)
- Use 16:9 aspect ratio (`prs.slide_width = Inches(13.333)`).
- One idea per slide.  Titles 32pt, body 18pt, source footer 10pt.
- House palette: ink `#0A0A0A`, accent `#0F62FE`, muted `#6B7280`,
  surface `#F5F5F5`.  No clipart.  No gradients unless requested.
- Every slide has a source footer (small, bottom-right).
- Recipe:
  ```python
  from pptx import Presentation
  from pptx.util import Inches, Pt
  prs = Presentation()
  prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
  s = prs.slides.add_slide(prs.slide_layouts[5])
  s.shapes.title.text = 'Title'
  tx = s.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.3), Inches(5))
  tf = tx.text_frame; tf.text = 'Body text'
  tf.paragraphs[0].runs[0].font.size = Pt(18)
  prs.save('/workspace/out.pptx')
  ```

## PDF from Markdown
- Preferred path: write the MD and let the UI convert.  The delivered
  MD file is already a first-class artefact.
- If the user insists on a server-side PDF, use `reportlab` (available
  in the sandbox).  Do NOT depend on pandoc, weasyprint, or
  wkhtmltopdf \u2014 not installed.
- Recipe (reportlab, simple):
  ```python
  from reportlab.lib.pagesizes import LETTER
  from reportlab.pdfgen import canvas
  c = canvas.Canvas('/workspace/out.pdf', pagesize=LETTER)
  c.setFont('Helvetica-Bold', 18); c.drawString(72, 720, 'Title')
  c.setFont('Helvetica', 11)
  y = 690
  for line in open('/workspace/out.md'):
      c.drawString(72, y, line.rstrip()[:95]); y -= 14
      if y < 72: c.showPage(); y = 720
  c.save()
  ```

## DOCX deliverables (python-docx)
- Not in the sandbox by default.  Attempt `pip install python-docx` in
  code_exec; if that fails, fall back to MD and tell the user.

## Final check before `deliver`
- Open the file you wrote; verify size > 0 and first 200 bytes look
  sane.  For XLSX/PPTX, re-open with openpyxl/python-pptx and count
  sheets/slides.  For PDF, verify header is `%PDF-`.
- If any check fails, fix before calling deliver.
""",
)


ALL_SKILLS: tuple[Skill, ...] = (
    CODING, WEB_DEV, DATA_ANALYSIS, BROWSER_AUTOMATION, DEVOPS, ML,
    CONTENT, RESEARCH, FINANCE, SECURITY,
    # v3 additions.
    SOCIAL_OPS, SALES_OPS, RECONCILIATION, FINANCE_MODELING, DOCUMENT_REVIEW,
    DELIVERABLE_QUALITY,
)


# ---------------------------------------------------------------------------
# Matcher
# ---------------------------------------------------------------------------


_WORD_RE = re.compile(r"[A-Za-z][A-Za-z0-9+#.-]*")


def _tokenise(text: str) -> set[str]:
    return {m.group(0).lower() for m in _WORD_RE.finditer(text or "")}


def score_skill(skill: Skill, goal_tokens: set[str], goal_text: str) -> int:
    """Score = single-word keyword hits + 2*multi-word keyword hits."""
    s = 0
    lower = goal_text.lower()
    for kw in skill.keywords:
        if " " in kw or "." in kw or "-" in kw or "+" in kw:
            if kw in lower:
                s += 2
        else:
            if kw in goal_tokens:
                s += 1
    return s


def select_skills(goal: str, user_instructions: str | None = None, top_k: int = 2) -> list[Skill]:
    """Return the top-k matching skills for a goal.  Empty list if nothing matches."""
    text = (goal or "") + "\n" + (user_instructions or "")
    tokens = _tokenise(text)
    scored = [(score_skill(s, tokens, text), s) for s in ALL_SKILLS]
    scored = [(sc, sk) for sc, sk in scored if sc > 0]
    scored.sort(key=lambda t: t[0], reverse=True)
    return [sk for _, sk in scored[:top_k]]


def render_skill_block(skills: list[Skill]) -> str:
    """Render a skills block ready to append to a system prompt."""
    if not skills:
        return ""
    parts = ["", "---", "DOMAIN SKILLS (apply these to the plan):"]
    for s in skills:
        parts.append("")
        parts.append(s.guidance.strip())
    return "\n".join(parts) + "\n"
