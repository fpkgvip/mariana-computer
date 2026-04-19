"""Document generation tools: PPTX, XLSX, DOCX.

These helpers produce Office-format files from structured data dicts.  They
are designed to be called from the orchestrator whenever the AI produces
structured output suitable for a deliverable document.

Dependencies: python-pptx, openpyxl, python-docx (listed in requirements.txt).
"""

from __future__ import annotations

import asyncio
from pathlib import Path

import structlog

logger = structlog.get_logger(__name__)


def _ensure_within_data_root(output_path: Path, data_root: Path | None) -> Path:
    """H-04 fix: resolve *output_path* and ensure it stays within *data_root*."""
    if data_root is None:
        return Path(output_path)
    resolved_root = Path(data_root).resolve()
    resolved_out = Path(output_path).resolve()
    if not resolved_out.is_relative_to(resolved_root):
        raise ValueError(
            f"Refusing to write document outside data_root: {output_path!r}"
        )
    return resolved_out


async def generate_pptx(
    title: str,
    slides_data: list[dict[str, object]],
    output_path: Path,
    data_root: Path | None = None,
) -> Path:
    """Generate a PowerPoint presentation.

    Parameters
    ----------
    title:
        Presentation title (shown on the first slide).
    slides_data:
        List of slide dicts, each with ``"title"`` and ``"bullets"`` keys.
    output_path:
        Destination file path.

    Returns
    -------
    Path
        The *output_path* after the file has been written.
    """
    def _build() -> Path:
        from pptx import Presentation  # noqa: PLC0415
        from pptx.util import Inches, Pt  # noqa: PLC0415

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # Title slide
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = title

        for slide_data in slides_data:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = str(slide_data.get("title", ""))
            if 1 in slide.shapes.placeholders:
                body = slide.shapes.placeholders[1]
                tf = body.text_frame
                for point in slide_data.get("bullets", []):
                    p = tf.add_paragraph()
                    p.text = str(point)
                    p.font.size = Pt(18)

        safe_out = _ensure_within_data_root(output_path, data_root)
        safe_out.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(safe_out))
        return safe_out

    result = await asyncio.to_thread(_build)
    logger.info("pptx_generated", title=title, slides=len(slides_data), path=str(output_path))
    return result


async def generate_xlsx(
    title: str,
    sheets_data: dict[str, list[list[object]]],
    output_path: Path,
    data_root: Path | None = None,
) -> Path:
    """Generate an Excel workbook.

    Parameters
    ----------
    title:
        Workbook title (used for logging; the first sheet is named from
        *sheets_data* keys).
    sheets_data:
        Mapping of sheet name -> list of rows (each row is a list of cells).
    output_path:
        Destination file path.

    Returns
    -------
    Path
        The *output_path* after the file has been written.
    """
    def _build() -> Path:
        from openpyxl import Workbook  # noqa: PLC0415

        wb = Workbook()
        first = True
        for sheet_name, rows in sheets_data.items():
            ws = wb.active if first else wb.create_sheet(sheet_name)
            if first:
                ws.title = sheet_name
                first = False
            for row in rows:
                ws.append(row)

        safe_out = _ensure_within_data_root(output_path, data_root)
        safe_out.parent.mkdir(parents=True, exist_ok=True)
        wb.save(str(safe_out))
        return safe_out

    result = await asyncio.to_thread(_build)
    logger.info("xlsx_generated", title=title, sheets=len(sheets_data), path=str(output_path))
    return result


async def generate_docx(
    title: str,
    sections: list[dict[str, object]],
    output_path: Path,
    data_root: Path | None = None,
) -> Path:
    """Generate a Word document.

    Parameters
    ----------
    title:
        Document title (inserted as a Heading 0).
    sections:
        List of section dicts.  Supported keys:

        - ``"heading"`` (str): Section heading text.
        - ``"level"`` (int): Heading level (1–4, default 1).
        - ``"text"`` (str): Body paragraph text.
        - ``"bullets"`` (list[str]): Bullet point list.
    output_path:
        Destination file path.

    Returns
    -------
    Path
        The *output_path* after the file has been written.
    """
    def _build() -> Path:
        from docx import Document  # noqa: PLC0415

        doc = Document()
        doc.add_heading(title, 0)

        for section in sections:
            heading = section.get("heading")
            if heading:
                level = int(section.get("level", 1))
                doc.add_heading(str(heading), level=level)
            text = section.get("text")
            if text:
                doc.add_paragraph(str(text))
            bullets = section.get("bullets")
            if bullets:
                for bullet in bullets:
                    doc.add_paragraph(str(bullet), style="List Bullet")

        safe_out = _ensure_within_data_root(output_path, data_root)
        safe_out.parent.mkdir(parents=True, exist_ok=True)
        doc.save(str(safe_out))
        return safe_out

    result = await asyncio.to_thread(_build)
    logger.info("docx_generated", title=title, sections=len(sections), path=str(output_path))
    return result
